import os
import subprocess
import shutil
from flask import send_file, after_this_request
from pptx import Presentation
import tempfile
import string
import random
import zipfile
from pathlib import Path
from pptx.util import Inches
from pdf2image import convert_from_path
import glob
from PIL import Image
from io import BytesIO


def pdf_to_pptx(file_storage):
    # Step 1: Create a random directory to store the files
    output_dir = os.path.join("/tmp", generate_random_string(10))  # Update the output directory

    os.makedirs(output_dir)

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp:
        temp_path = temp.name
        file_storage.save(temp_path)

    subprocess.run(["gs", "-dNOPAUSE", "-dBATCH", "-sDEVICE=jpeg", "-r150", f"-sOutputFile={output_dir}/%d.jpg", temp_path])
    os.remove(temp_path)

    # Step 2: Read the generated images and insert them into the PowerPoint
    prs = Presentation()

    jpg_files = sorted(os.listdir(output_dir), key=lambda x: int(''.join(filter(str.isdigit, os.path.splitext(x)[0]))))
    for jpg_file in jpg_files:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.add_picture(os.path.join(output_dir, jpg_file), 0, 0)

    # Step 3: Save the PowerPoint file
    pptx_file = os.path.join(output_dir, "converted_file.pptx")
    prs.save(pptx_file)

    # Step 4: Delete the directory after the current request
    @after_this_request
    def delete_directory(response):
        shutil.rmtree(output_dir)
        return response

    # Step 5: Return the PowerPoint file for download
    return send_file(pptx_file, as_attachment=True, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

def generate_random_string(length):
    letters = string.ascii_lowercase
    return ''.join(random.choice(letters) for _ in range(length))





"""
    please update this funciton to adjust the powerpoint slide to be == the size of each image.
"""

def pdf_to_pptx_multiple(pdf_files):
    output_dir = tempfile.mkdtemp()

    for pdf_file in pdf_files:
        pdf_filename = Path(pdf_file.filename).stem
        
        # Extract PDF pages as images using pdf2image
        with tempfile.NamedTemporaryFile(suffix=".pdf") as temp_pdf:
            pdf_file.save(temp_pdf.name)
            images = convert_from_path(temp_pdf.name, dpi=300) # This returns a list of PIL Image objects

        # Create PPTX file and add images
        prs = Presentation() 
        for img in images:
            # Change the slide size to match the image size in inches
            prs.slide_width = Inches(img.width / 100)
            prs.slide_height = Inches(img.height / 100)

            slide = prs.slides.add_slide(prs.slide_layouts[1]) 
            # Save and crop the image to a BytesIO object and add it to the slide
            temp_img = BytesIO() # This creates a BytesIO object
            img.save(temp_img, format="JPEG") # This saves the image to the BytesIO object
            # Calculate the new width and height based on the aspect ratio
            new_width = img.height * 16 / 9
            new_height = img.width * 9 / 16

            # Calculate the coordinates for cropping
            left = (img.width - new_width) / 2
            top = (img.height - new_height) / 2
            right = (img.width + new_width) / 2
            bottom = (img.height + new_height) / 2

            # Crop the image using PIL
            img = Image.open(temp_img)
            img = img.crop((left, top, right, bottom))
            img.save(temp_img, format="JPEG") # This saves the cropped image to the same BytesIO object

            slide.shapes.add_picture(temp_img, 0, 0) # This adds the image from the BytesIO object to the slide

        pptx_path = os.path.join(output_dir, f"{pdf_filename}.pptx")
        prs.save(pptx_path)

    # Zip all PPTX files
    zip_path = os.path.join(output_dir, "converted_files.zip")
    with zipfile.ZipFile(zip_path, "w") as zip_f:
        for pptx_file in glob.glob(os.path.join(output_dir, "*.pptx")):
            zip_f.write(pptx_file, os.path.basename(pptx_file))


    @after_this_request
    def cleanup(response):
        shutil.rmtree(output_dir) # This will delete the output directory
        return response

    # Send the zip file as a response
    return send_file(zip_path, mimetype="application/zip", download_name="converted_files.zip", as_attachment=True)
